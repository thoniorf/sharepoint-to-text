"""
PPTX content extractor using python-pptx library.
"""

import io
import logging
import zipfile
from datetime import datetime
from typing import Any, Generator, List, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from sharepoint2text.extractors.data_types import (
    PPTXComment,
    PptxContent,
    PPTXFormula,
    PPTXImage,
    PptxMetadata,
    PPTXSlide,
)
from sharepoint2text.extractors.docx_extractor import _DocxFullTextExtractor

logger = logging.getLogger(__name__)


def _extract_slide_comments(
    file_like: io.BytesIO, slide_number: int
) -> List[PPTXComment]:
    """Extract comments for a specific slide from the PPTX archive."""
    comments = []
    file_like.seek(0)

    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"

    try:
        with zipfile.ZipFile(file_like, "r") as z:
            comment_file = f"ppt/comments/comment{slide_number}.xml"
            if comment_file not in z.namelist():
                return comments

            with z.open(comment_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()

            for cm in root.findall(f".//{{{p_ns}}}cm"):
                author_id = cm.get("authorId", "")
                text_elem = cm.find(f"{{{p_ns}}}text")
                text = text_elem.text if text_elem is not None else ""
                dt = cm.get("dt", "")
                comments.append(
                    PPTXComment(
                        author=author_id,
                        text=text or "",
                        date=dt,
                    )
                )
    except Exception as e:
        logger.debug(f"Failed to extract comments for slide {slide_number}: {e}")

    return comments


def _extract_formulas_from_shape(shape) -> List[Tuple[str, bool]]:
    """Extract formulas from a shape's XML, returning list of (latex, is_display) tuples."""
    formulas = []
    m_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/math}"

    try:
        # Access the shape's XML element
        shape_elem = shape._element

        # Find all oMathPara (display equations) and oMath (inline equations)
        for elem in shape_elem.iter():
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

            if tag == "oMathPara":
                omath = elem.find(f"{m_ns}oMath")
                if omath is not None:
                    latex = _DocxFullTextExtractor.omml_to_latex(omath)
                    if latex.strip():
                        formulas.append((latex, True))

            elif tag == "oMath":
                # Check if parent is oMathPara - if so, skip (already handled)
                parent = None
                try:
                    parent = elem.getparent()
                except Exception:
                    pass
                if parent is not None:
                    parent_tag = (
                        parent.tag.split("}")[-1] if "}" in parent.tag else parent.tag
                    )
                    if parent_tag == "oMathPara":
                        continue
                latex = _DocxFullTextExtractor.omml_to_latex(elem)
                if latex.strip():
                    formulas.append((latex, False))
    except Exception as e:
        logger.debug(f"Failed to extract formulas from shape: {e}")

    return formulas


def _get_shape_position(shape) -> Tuple[int, int]:
    """Get the position (left, top) of a shape for ordering purposes."""
    try:
        left = shape.left if shape.left is not None else 0
        top = shape.top if shape.top is not None else 0
        return (top, left)  # Sort by top first, then left
    except Exception:
        return (0, 0)


def _dt_to_iso(dt: datetime | None) -> str:
    return dt.isoformat() if dt else ""


def read_pptx(
    file_like: io.BytesIO, path: str | None = None
) -> Generator[PptxContent, Any, None]:
    """
    Extract all relevant content from a PPTX file.

    Args:
        file_like: A BytesIO object containing the PPTX file data.
        path: Optional file path to populate file metadata fields.

    Yields:
        MicrosoftPptxContent dataclass with all extracted content.
    """
    logger.debug("Reading pptx")
    file_like.seek(0)
    prs = Presentation(file_like)

    cp = prs.core_properties
    metadata = PptxMetadata(
        title=cp.title or "",
        subject=cp.subject or "",
        author=cp.author or "",
        last_modified_by=cp.last_modified_by or "",
        created=_dt_to_iso(cp.created),
        modified=_dt_to_iso(cp.modified),
        keywords=cp.keywords or "",
        comments=cp.comments or "",
        category=cp.category or "",
        revision=cp.revision,
    )

    slides_result: List[PPTXSlide] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        logger.debug(f"Processing slide [{slide_index}]")

        slide_title = ""
        slide_footer = ""
        content_placeholders: List[str] = []
        other_textboxes: List[str] = []
        images: List[PPTXImage] = []
        formulas: List[PPTXFormula] = []

        image_counter = 0

        # Collect all content items with their positions for ordering
        # Each item: (position, content_type, content_text)
        ordered_content: List[Tuple[Tuple[int, int], str, str]] = []

        # Sort shapes by position (top to bottom, left to right)
        sorted_shapes = sorted(slide.shapes, key=_get_shape_position)

        for shape in sorted_shapes:
            position = _get_shape_position(shape)

            # ---------------------------
            # Image extraction with caption
            # ---------------------------
            if shape.shape_type == shape.shape_type.PICTURE:
                try:
                    image = shape.image
                    image_counter += 1

                    # Extract alt text / description
                    caption = ""
                    try:
                        # Alt text is in the shape's description property
                        if hasattr(shape, "name") and shape.name:
                            caption = shape.name
                        # Check for description in non-visual properties
                        desc_cNvPr = shape._element.find(
                            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr"
                        )
                        if desc_cNvPr is not None:
                            descr = desc_cNvPr.get("descr", "")
                            if descr:
                                caption = descr
                    except Exception as e:
                        logger.debug(f"Failed to extract image caption: {e}")

                    images.append(
                        PPTXImage(
                            image_index=image_counter,
                            filename=image.filename,
                            content_type=image.content_type,
                            size_bytes=len(image.blob),
                            blob=image.blob,
                            caption=caption,
                        )
                    )

                    # Add caption to ordered content if present
                    if caption:
                        ordered_content.append(
                            (position, "image_caption", f"[Image: {caption}]")
                        )

                except Exception as e:
                    logger.error(e)
                    logger.exception(f"Failed to extract image on slide {slide_index}")
                continue

            # ---------------------------
            # Formula extraction
            # ---------------------------
            shape_formulas = _extract_formulas_from_shape(shape)
            for latex, is_display in shape_formulas:
                formula = PPTXFormula(latex=latex, is_display=is_display)
                formulas.append(formula)
                if is_display:
                    ordered_content.append((position, "formula", f"$${latex}$$"))
                else:
                    ordered_content.append((position, "formula", f"${latex}$"))

            # ---------------------------
            # Text extraction
            # ---------------------------
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            if shape.is_placeholder:
                ptype = shape.placeholder_format.type

                if ptype in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                    PP_PLACEHOLDER.VERTICAL_TITLE,
                ):
                    slide_title = text
                    ordered_content.append((position, "title", text))

                elif ptype == PP_PLACEHOLDER.FOOTER:
                    slide_footer = text
                    # Footer is typically not included in main text

                elif ptype in (
                    PP_PLACEHOLDER.BODY,
                    PP_PLACEHOLDER.SUBTITLE,
                    PP_PLACEHOLDER.OBJECT,
                    PP_PLACEHOLDER.VERTICAL_BODY,
                    PP_PLACEHOLDER.VERTICAL_OBJECT,
                    PP_PLACEHOLDER.TABLE,
                ):
                    content_placeholders.append(text)
                    ordered_content.append((position, "content", text))

                else:
                    other_textboxes.append(text)
                    ordered_content.append((position, "other", text))
            else:
                other_textboxes.append(text)
                ordered_content.append((position, "other", text))

        # ---------------------------
        # Comment extraction
        # ---------------------------
        comments = _extract_slide_comments(file_like, slide_index)
        # Add comments at the end of the slide content
        for comment in comments:
            ordered_content.append(
                (
                    (999999, 999999),
                    "comment",
                    f"[Comment: {comment.author}@{comment.date}: {comment.text}]",
                )
            )

        # Build slide text from ordered content
        # Sort by position (already sorted but ensure stability)
        ordered_content.sort(key=lambda x: x[0])
        slide_text_parts = [item[2] for item in ordered_content]
        slide_text = "\n".join(slide_text_parts)

        # Build base text (without formulas, comments, image captions)
        base_content_types = {"title", "content", "other"}
        base_text_parts = [
            item[2] for item in ordered_content if item[1] in base_content_types
        ]
        base_text = "\n".join(base_text_parts)

        slides_result.append(
            PPTXSlide(
                slide_number=slide_index,
                title=slide_title,
                footer=slide_footer,
                content_placeholders=content_placeholders,
                other_textboxes=other_textboxes,
                images=images,
                formulas=formulas,
                comments=comments,
                text=slide_text,
                base_text=base_text,
            )
        )

    metadata.populate_from_path(path)
    yield PptxContent(metadata=metadata, slides=slides_result)
