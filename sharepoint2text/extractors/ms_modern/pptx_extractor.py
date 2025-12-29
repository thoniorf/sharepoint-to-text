"""
PPTX Presentation Extractor
===========================

Extracts text content, metadata, and structure from Microsoft PowerPoint .pptx
files (Office Open XML format, PowerPoint 2007 and later).

This module uses the python-pptx library for high-level presentation access and
direct XML parsing for specialized content (comments, embedded formulas).

File Format Background
----------------------
The .pptx format is a ZIP archive containing XML files following the Office
Open XML (OOXML) standard. Key components:

    ppt/presentation.xml: Presentation-level properties
    ppt/slides/slide1.xml, slide2.xml, ...: Individual slide content
    ppt/slideLayouts/: Slide layout templates
    ppt/slideMasters/: Master slide definitions
    ppt/comments/comment1.xml, ...: Per-slide comments
    ppt/media/: Embedded images and media
    docProps/core.xml: Metadata (title, author, dates)

XML Namespaces:
    - p: http://schemas.openxmlformats.org/presentationml/2006/main
    - a: http://schemas.openxmlformats.org/drawingml/2006/main
    - m: http://schemas.openxmlformats.org/officeDocument/2006/math
    - r: http://schemas.openxmlformats.org/officeDocument/2006/relationships

Dependencies
------------
python-pptx: https://github.com/scanny/python-pptx
    pip install python-pptx

    Provides:
    - Slide enumeration and access
    - Shape parsing (text boxes, images, placeholders)
    - Placeholder type detection
    - Image extraction
    - Core properties (metadata)

Math Formula Handling
---------------------
PowerPoint can contain math formulas in OMML format (same as Word).
This module reuses the OMML-to-LaTeX converter from docx_extractor
to extract formulas as LaTeX notation.

Shape Types and Placeholders
----------------------------
PowerPoint shapes are categorized by type and placeholder function:

Placeholder Types (PP_PLACEHOLDER):
    - TITLE, CENTER_TITLE, VERTICAL_TITLE: Slide titles
    - BODY, SUBTITLE: Main content areas
    - FOOTER: Footer text
    - OBJECT, TABLE: Content containers

Text Ordering:
    Shapes are sorted by position (top-to-bottom, left-to-right) to
    maintain a logical reading order in the extracted text.

Extracted Content
-----------------
Per-slide content includes:
    - title: Slide title text
    - content_placeholders: Body text from content areas
    - other_textboxes: Text from non-placeholder shapes
    - images: Embedded images with metadata and binary data
    - formulas: Math formulas as LaTeX
    - comments: Slide comments with author and date
    - text: Complete slide text in reading order
    - base_text: Text without formulas/comments/captions

Known Limitations
-----------------
- SmartArt text extraction may be incomplete
- Chart data/labels are not extracted as text
- Grouped shapes may not extract all nested text
- Speaker notes are not currently extracted
- Audio/video content is not extracted
- Password-protected files are not supported
- Very large presentations may use significant memory

Usage
-----
    >>> import io
    >>> from sharepoint2text.extractors.ms_modern.pptx_extractor import read_pptx
    >>>
    >>> with open("slides.pptx", "rb") as f:
    ...     for ppt in read_pptx(io.BytesIO(f.read()), path="slides.pptx"):
    ...         print(f"Title: {ppt.metadata.title}")
    ...         for slide in ppt.slides:
    ...             print(f"Slide {slide.slide_number}: {slide.title}")
    ...             print(slide.text)

See Also
--------
- python-pptx documentation: https://python-pptx.readthedocs.io/
- ppt_extractor: For legacy .ppt format

Maintenance Notes
-----------------
- Shape position sorting ensures consistent text order
- Comment extraction requires direct XML parsing (not in python-pptx API)
- Formula extraction reuses docx_extractor's OMML converter
- Image alt text is extracted from cNvPr descr attribute
"""

import io
import logging
import zipfile
from datetime import datetime
from typing import Any, Generator, List, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from sharepoint2text.extractors.data_types import (
    PPTXComment,
    PptxContent,
    PPTXFormula,
    PPTXImage,
    PptxMetadata,
    PPTXSlide,
)
from sharepoint2text.extractors.ms_modern.docx_extractor import _DocxFullTextExtractor

logger = logging.getLogger(__name__)


def _extract_slide_comments(
    file_like: io.BytesIO, slide_number: int
) -> List[PPTXComment]:
    """
    Extract comments for a specific slide by parsing the comments XML.

    Comments in PPTX files are stored in separate XML files, one per slide
    that has comments (ppt/comments/comment{n}.xml).

    Args:
        file_like: BytesIO containing the PPTX file.
        slide_number: 1-based slide number to extract comments for.

    Returns:
        List of PPTXComment objects with author, text, and date fields.
        Returns empty list if no comment file exists for the slide.

    Notes:
        - Author is stored as authorId (numeric), not the actual name
        - Date is in ISO format from the XML
    """
    comments = []
    file_like.seek(0)

    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"

    try:
        with zipfile.ZipFile(file_like, "r") as z:
            comment_file = f"ppt/comments/comment{slide_number}.xml"
            if comment_file not in z.namelist():
                return comments

            with z.open(comment_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()

            for cm in root.findall(f".//{{{p_ns}}}cm"):
                author_id = cm.get("authorId", "")
                text_elem = cm.find(f"{{{p_ns}}}text")
                text = text_elem.text if text_elem is not None else ""
                dt = cm.get("dt", "")
                comments.append(
                    PPTXComment(
                        author=author_id,
                        text=text or "",
                        date=dt,
                    )
                )
    except Exception as e:
        logger.debug(f"Failed to extract comments for slide {slide_number}: {e}")

    return comments


def _extract_formulas_from_shape(shape) -> List[Tuple[str, bool]]:
    """
    Extract mathematical formulas from a shape's XML content.

    Searches the shape's XML element tree for OMML math elements and
    converts them to LaTeX using the shared OMML converter.

    Args:
        shape: python-pptx Shape object.

    Returns:
        List of tuples (latex_string, is_display) where:
        - latex_string: LaTeX representation of the formula
        - is_display: True for display equations (oMathPara), False for inline

    Notes:
        - Uses _DocxFullTextExtractor.omml_to_latex for conversion
        - Handles both display (oMathPara) and inline (oMath) equations
        - Skips inline equations that are children of display equations
    """
    formulas = []
    m_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/math}"

    try:
        # Access the shape's XML element
        shape_elem = shape._element

        # Find all oMathPara (display equations) and oMath (inline equations)
        for elem in shape_elem.iter():
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

            if tag == "oMathPara":
                omath = elem.find(f"{m_ns}oMath")
                if omath is not None:
                    latex = _DocxFullTextExtractor.omml_to_latex(omath)
                    if latex.strip():
                        formulas.append((latex, True))

            elif tag == "oMath":
                # Check if parent is oMathPara - if so, skip (already handled)
                parent = None
                try:
                    parent = elem.getparent()
                except Exception:
                    pass
                if parent is not None:
                    parent_tag = (
                        parent.tag.split("}")[-1] if "}" in parent.tag else parent.tag
                    )
                    if parent_tag == "oMathPara":
                        continue
                latex = _DocxFullTextExtractor.omml_to_latex(elem)
                if latex.strip():
                    formulas.append((latex, False))
    except Exception as e:
        logger.debug(f"Failed to extract formulas from shape: {e}")

    return formulas


def _get_shape_position(shape) -> Tuple[int, int]:
    """
    Get the position of a shape for sorting purposes.

    Returns position as (top, left) tuple so shapes can be sorted
    in reading order (top-to-bottom, then left-to-right).

    Args:
        shape: python-pptx Shape object.

    Returns:
        Tuple of (top, left) coordinates in EMUs. Returns (0, 0) if
        position cannot be determined.
    """
    try:
        left = shape.left if shape.left is not None else 0
        top = shape.top if shape.top is not None else 0
        return (top, left)  # Sort by top first, then left
    except Exception:
        return (0, 0)


def _dt_to_iso(dt: datetime | None) -> str:
    """Convert datetime to ISO format string, or empty string if None."""
    return dt.isoformat() if dt else ""


def read_pptx(
    file_like: io.BytesIO, path: str | None = None
) -> Generator[PptxContent, Any, None]:
    """
    Extract all relevant content from a PowerPoint .pptx file.

    Primary entry point for PPTX file extraction. Iterates through all slides,
    extracting text, images, formulas, and comments while maintaining shape
    ordering for consistent text output.

    This function uses a generator pattern for API consistency with other
    extractors, even though PPTX files contain exactly one presentation.

    Args:
        file_like: BytesIO object containing the complete PPTX file data.
            The stream position is reset to the beginning before reading.
        path: Optional filesystem path to the source file. If provided,
            populates file metadata (filename, extension, folder) in the
            returned PptxContent.metadata.

    Yields:
        PptxContent: Single PptxContent object containing:
            - metadata: PptxMetadata with title, author, dates, revision
            - slides: List of PPTXSlide objects, each containing:
                - slide_number: 1-based slide index
                - title: Slide title text
                - content_placeholders: Body text from content areas
                - other_textboxes: Text from non-placeholder shapes
                - images: List of PPTXImage with binary data
                - formulas: List of PPTXFormula as LaTeX
                - comments: List of PPTXComment
                - text: Complete slide text with formulas and comments
                - base_text: Text without formulas/comments/captions

    Processing Details:
        - Shapes are sorted by position (top-to-bottom, left-to-right)
        - Title placeholders are extracted separately from body content
        - Images include alt text/captions when available
        - Formulas are converted to LaTeX notation
        - Comments are appended at the end of slide content

    Example:
        >>> import io
        >>> with open("presentation.pptx", "rb") as f:
        ...     data = io.BytesIO(f.read())
        ...     for ppt in read_pptx(data, path="presentation.pptx"):
        ...         print(f"Title: {ppt.metadata.title}")
        ...         print(f"Slides: {len(ppt.slides)}")
        ...         for slide in ppt.slides:
        ...             print(f"  Slide {slide.slide_number}: {slide.title}")
        ...             print(f"    Images: {len(slide.images)}")

    Performance Notes:
        - Images are loaded into memory as binary blobs
        - Large presentations with many images may use significant memory
        - Comments require a separate ZIP file read per slide
    """
    logger.debug("Reading pptx")
    file_like.seek(0)
    prs = Presentation(file_like)

    cp = prs.core_properties
    metadata = PptxMetadata(
        title=cp.title or "",
        subject=cp.subject or "",
        author=cp.author or "",
        last_modified_by=cp.last_modified_by or "",
        created=_dt_to_iso(cp.created),
        modified=_dt_to_iso(cp.modified),
        keywords=cp.keywords or "",
        comments=cp.comments or "",
        category=cp.category or "",
        revision=cp.revision,
    )

    slides_result: List[PPTXSlide] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        logger.debug(f"Processing slide [{slide_index}]")

        slide_title = ""
        slide_footer = ""
        content_placeholders: List[str] = []
        other_textboxes: List[str] = []
        images: List[PPTXImage] = []
        formulas: List[PPTXFormula] = []

        image_counter = 0

        # Collect all content items with their positions for ordering
        # Each item: (position, content_type, content_text)
        ordered_content: List[Tuple[Tuple[int, int], str, str]] = []

        # Sort shapes by position (top to bottom, left to right)
        sorted_shapes = sorted(slide.shapes, key=_get_shape_position)

        for shape in sorted_shapes:
            position = _get_shape_position(shape)

            # ---------------------------
            # Image extraction with caption
            # ---------------------------
            if shape.shape_type == shape.shape_type.PICTURE:
                try:
                    image = shape.image
                    image_counter += 1

                    # Extract alt text / description
                    caption = ""
                    try:
                        # Alt text is in the shape's description property
                        if hasattr(shape, "name") and shape.name:
                            caption = shape.name
                        # Check for description in non-visual properties
                        desc_cNvPr = shape._element.find(
                            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr"
                        )
                        if desc_cNvPr is not None:
                            descr = desc_cNvPr.get("descr", "")
                            if descr:
                                caption = descr
                    except Exception as e:
                        logger.debug(f"Failed to extract image caption: {e}")

                    images.append(
                        PPTXImage(
                            image_index=image_counter,
                            filename=image.filename,
                            content_type=image.content_type,
                            size_bytes=len(image.blob),
                            blob=image.blob,
                            caption=caption,
                        )
                    )

                    # Add caption to ordered content if present
                    if caption:
                        ordered_content.append(
                            (position, "image_caption", f"[Image: {caption}]")
                        )

                except Exception as e:
                    logger.error(e)
                    logger.exception(f"Failed to extract image on slide {slide_index}")
                continue

            # ---------------------------
            # Formula extraction
            # ---------------------------
            shape_formulas = _extract_formulas_from_shape(shape)
            for latex, is_display in shape_formulas:
                formula = PPTXFormula(latex=latex, is_display=is_display)
                formulas.append(formula)
                if is_display:
                    ordered_content.append((position, "formula", f"$${latex}$$"))
                else:
                    ordered_content.append((position, "formula", f"${latex}$"))

            # ---------------------------
            # Text extraction
            # ---------------------------
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            if shape.is_placeholder:
                ptype = shape.placeholder_format.type

                if ptype in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                    PP_PLACEHOLDER.VERTICAL_TITLE,
                ):
                    slide_title = text
                    ordered_content.append((position, "title", text))

                elif ptype == PP_PLACEHOLDER.FOOTER:
                    slide_footer = text
                    # Footer is typically not included in main text

                elif ptype in (
                    PP_PLACEHOLDER.BODY,
                    PP_PLACEHOLDER.SUBTITLE,
                    PP_PLACEHOLDER.OBJECT,
                    PP_PLACEHOLDER.VERTICAL_BODY,
                    PP_PLACEHOLDER.VERTICAL_OBJECT,
                    PP_PLACEHOLDER.TABLE,
                ):
                    content_placeholders.append(text)
                    ordered_content.append((position, "content", text))

                else:
                    other_textboxes.append(text)
                    ordered_content.append((position, "other", text))
            else:
                other_textboxes.append(text)
                ordered_content.append((position, "other", text))

        # ---------------------------
        # Comment extraction
        # ---------------------------
        comments = _extract_slide_comments(file_like, slide_index)
        # Add comments at the end of the slide content
        for comment in comments:
            ordered_content.append(
                (
                    (999999, 999999),
                    "comment",
                    f"[Comment: {comment.author}@{comment.date}: {comment.text}]",
                )
            )

        # Build slide text from ordered content
        # Sort by position (already sorted but ensure stability)
        ordered_content.sort(key=lambda x: x[0])
        slide_text_parts = [item[2] for item in ordered_content]
        slide_text = "\n".join(slide_text_parts)

        # Build base text (without formulas, comments, image captions)
        base_content_types = {"title", "content", "other"}
        base_text_parts = [
            item[2] for item in ordered_content if item[1] in base_content_types
        ]
        base_text = "\n".join(base_text_parts)

        slides_result.append(
            PPTXSlide(
                slide_number=slide_index,
                title=slide_title,
                footer=slide_footer,
                content_placeholders=content_placeholders,
                other_textboxes=other_textboxes,
                images=images,
                formulas=formulas,
                comments=comments,
                text=slide_text,
                base_text=base_text,
            )
        )

    metadata.populate_from_path(path)
    yield PptxContent(metadata=metadata, slides=slides_result)
